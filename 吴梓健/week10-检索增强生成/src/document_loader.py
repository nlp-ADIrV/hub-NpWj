"""
文档解析模块 — 支持 PDF 和 PPT 两种格式

设计思路：
  1. 统一的 ParsedBlock 数据结构，无论来源是 PDF 还是 PPT
  2. PDF 解析：pdfplumber（表格）+ PyMuPDF/fitz（文字+字体信息）
  3. PPT 解析：python-pptx（幻灯片文字 + 表格 + 备注）
  4. 自动识别文件后缀，分发给对应的解析器
  5. 保留页码/幻灯片号、章节路径等元信息，供后续溯源

输出格式（与原项目一致）：
  每个 ParsedBlock 包含：
    block_type   : "text" | "table" | "title"
    content      : 文字内容（表格转为 Markdown）
    page_num     : 页码或幻灯片编号
    section_path : 章节路径栈
    is_ocr       : 是否经过 OCR
    source_file  : 来源文件名

使用方式：
  from document_loader import load_document, load_all_documents
  blocks = load_document(Path("data/raw_docs/report.pdf"))
  # 或批量加载
  all_blocks = load_all_documents()

  # 命令行直接运行
  python src/document_loader.py
"""

import re
import json
import logging
from pathlib import Path
from dataclasses import dataclass, field, asdict
from typing import Optional

from config import settings

logger = logging.getLogger(__name__)


# ── 数据结构 ──────────────────────────────────────────────────────────────────

@dataclass
class ParsedBlock:
    """
    解析块 = 一段连续内容（文字段落 / 表格 / 标题）

    无论来自 PDF 还是 PPT，统一用这个结构表示，
    方便后续分块和向量化阶段统一处理。
    """
    block_type:   str            # "text" | "table" | "title"
    content:      str            # 文字内容（表格转为 markdown）
    page_num:     int            # PDF 页码 / PPT 幻灯片编号
    section_path: list[str]      # 章节路径，如 ["第一章 概述", "1.1 背景"]
    is_ocr:       bool = False
    source_file:  str = ""
    raw_table:    Optional[list] = field(default=None, repr=False)


# ── 工具函数 ──────────────────────────────────────────────────────────────────

# 章节标题模式
CHAPTER_PATTERNS = [
    re.compile(r"^第[一二三四五六七八九十百]+[章节]"),
    re.compile(r"^[一二三四五六七八九十]、"),
    re.compile(r"^\d+\.\s"),
    re.compile(r"^\d+\.\d+\s"),       # 1.1 标题
]

# 噪声行模式（页眉/页脚/页码）
NOISE_PATTERNS = [
    re.compile(r"^.{1,40}年度报告\s*$"),
    re.compile(r"^\d+\s*$"),
    re.compile(r"^—\s*\d+\s*—$"),
]


def is_noise_line(line: str) -> bool:
    """判断是否为页眉/页脚噪声行。"""
    line = line.strip()
    if len(line) < 2:
        return True
    return any(p.match(line) for p in NOISE_PATTERNS)


def is_title_line(line: str, fontsize: Optional[float] = None, is_bold: bool = False) -> bool:
    """判断一行是否为标题（有字体信息时优先用字体大小判断）。"""
    if fontsize and fontsize >= 14:
        return True
    if is_bold and len(line.strip()) < 50:
        return True
    return any(p.match(line.strip()) for p in CHAPTER_PATTERNS)


def table_to_markdown(table: list[list]) -> str:
    """将二维表格数据转为 Markdown 格式，方便 LLM 读取。"""
    if not table:
        return ""

    rows = []
    for row in table:
        cleaned = [str(cell or "").replace("\n", " ").strip() for cell in row]
        rows.append(cleaned)

    if not rows:
        return ""

    header = rows[0]
    lines = ["| " + " | ".join(header) + " |"]
    lines.append("| " + " | ".join(["---"] * len(header)) + " |")
    for row in rows[1:]:
        while len(row) < len(header):
            row.append("")
        lines.append("| " + " | ".join(row[:len(header)]) + " |")

    return "\n".join(lines)


# ── PDF 解析器 ────────────────────────────────────────────────────────────────

class PDFParser:
    """
    PDF 文档解析器。

    策略：
      - pdfplumber 提取表格（行列识别更准确）
      - PyMuPDF/fitz 提取带字体信息的文字（用于判断标题层级）
      - 每页维护章节路径栈
    """

    def __init__(self, pdf_path: Path):
        self.pdf_path = pdf_path
        self.blocks: list[ParsedBlock] = []
        self._section_stack: list[str] = []

    def _update_section(self, title: str):
        """维护章节栈，根据编号层级推断层次。"""
        if re.match(r"^第[一二三四五六七八九十]+章", title):
            self._section_stack = [title]
        elif re.match(r"^第[一二三四五六七八九十]+节", title):
            self._section_stack = self._section_stack[:1] + [title]
        elif re.match(r"^[一二三四五六七八九十]、", title):
            self._section_stack = self._section_stack[:2] + [title]
        elif re.match(r"^\d+\.\d+", title):
            self._section_stack = self._section_stack[:3] + [title]
        else:
            self._section_stack = self._section_stack[:3] + [title]

    def parse(self) -> list[ParsedBlock]:
        import pdfplumber
        import fitz  # PyMuPDF

        logger.info(f"解析 PDF: {self.pdf_path.name}")

        plumber_doc = pdfplumber.open(self.pdf_path)
        fitz_doc = fitz.open(str(self.pdf_path))

        for page_num in range(len(fitz_doc)):
            fitz_page = fitz_doc[page_num]
            plumb_page = plumber_doc.pages[page_num]

            # ── 提取表格 ──
            for table in plumb_page.extract_tables():
                if table:
                    md = table_to_markdown(table)
                    if md:
                        self.blocks.append(ParsedBlock(
                            block_type="table",
                            content=md,
                            page_num=page_num + 1,
                            section_path=list(self._section_stack),
                            source_file=self.pdf_path.name,
                            raw_table=table,
                        ))

            # ── 提取文字（带字体信息）──
            page_dict = fitz_page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
            current_para_lines = []

            for block in page_dict.get("blocks", []):
                if block.get("type") != 0:  # 0=文字，1=图片
                    continue

                for line in block.get("lines", []):
                    line_text = "".join(
                        span["text"] for span in line.get("spans", [])
                    ).strip()

                    if not line_text or is_noise_line(line_text):
                        continue

                    spans = line.get("spans", [])
                    fontsize = spans[0].get("size", 0) if spans else 0
                    is_bold = any("Bold" in span.get("font", "") for span in spans)

                    if is_title_line(line_text, fontsize, is_bold):
                        if current_para_lines:
                            self.blocks.append(ParsedBlock(
                                block_type="text",
                                content="\n".join(current_para_lines),
                                page_num=page_num + 1,
                                section_path=list(self._section_stack),
                                source_file=self.pdf_path.name,
                            ))
                            current_para_lines = []

                        self._update_section(line_text)
                        self.blocks.append(ParsedBlock(
                            block_type="title",
                            content=line_text,
                            page_num=page_num + 1,
                            section_path=list(self._section_stack),
                            source_file=self.pdf_path.name,
                        ))
                    else:
                        current_para_lines.append(line_text)

            if current_para_lines:
                self.blocks.append(ParsedBlock(
                    block_type="text",
                    content="\n".join(current_para_lines),
                    page_num=page_num + 1,
                    section_path=list(self._section_stack),
                    source_file=self.pdf_path.name,
                ))

        plumber_doc.close()
        fitz_doc.close()

        logger.info(f"  PDF 解析完成: {len(self.blocks)} 个块")
        return self.blocks


# ── PPT 解析器 ────────────────────────────────────────────────────────────────

class PPTParser:
    """
    PowerPoint 文档解析器。

    策略：
      - 遍历每张幻灯片，提取标题 + 正文文本框 + 表格 + 备注页
      - 幻灯片标题识别为 title 块（用于章节路径维护）
      - 表格转为 Markdown 格式
      - 备注页文本作为额外 text 块（PPT 演讲者备注常含重要信息）
    """

    def __init__(self, ppt_path: Path):
        self.ppt_path = ppt_path
        self.blocks: list[ParsedBlock] = []
        self._section_stack: list[str] = []

    def _extract_table(self, shape) -> str:
        """从 PPT 表格 shape 中提取 Markdown 表格。"""
        table = shape.table
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append(cells)
        return table_to_markdown(rows)

    def _iter_shapes(self, shapes):
        """
        递归遍历所有 shape，包括 GroupShape 内部嵌套的 shape。

        PPT 中表格、文本框常被放在 GroupShape 里，
        slide.shapes 只返回顶层，需要递归才能取到。
        """
        for shape in shapes:
            # 递归处理 GroupShape
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                yield from self._iter_shapes(shape.shapes)
            else:
                yield shape

    def parse(self) -> list[ParsedBlock]:
        from pptx import Presentation

        logger.info(f"解析 PPT: {self.ppt_path.name}")

        prs = Presentation(str(self.ppt_path))

        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            slide_title = ""

            # 收集本页所有 shape（递归展开 GroupShape）
            all_shapes = list(self._iter_shapes(slide.shapes))

            # ── 提取标题（通常在第一个 placeholder）──
            for shape in all_shapes:
                if shape.has_text_frame and shape.is_placeholder:
                    if shape.placeholder_format.idx == 0:  # 标题占位符
                        slide_title = shape.text_frame.text.strip()
                        if slide_title:
                            self._section_stack = [slide_title]
                            self.blocks.append(ParsedBlock(
                                block_type="title",
                                content=slide_title,
                                page_num=slide_num,
                                section_path=list(self._section_stack),
                                source_file=self.ppt_path.name,
                            ))
                        break

            # ── 提取正文文本框和表格 ──
            for shape in all_shapes:
                # 跳过已处理的标题
                if shape.is_placeholder and shape.placeholder_format.idx == 0:
                    continue

                # 表格
                if shape.has_table:
                    md = self._extract_table(shape)
                    if md:
                        self.blocks.append(ParsedBlock(
                            block_type="table",
                            content=md,
                            page_num=slide_num,
                            section_path=list(self._section_stack),
                            source_file=self.ppt_path.name,
                        ))

                # 文本框
                elif shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text and len(text) > 2:
                        # 判断是否是子标题
                        if is_title_line(text):
                            self._section_stack = self._section_stack[:1] + [text]
                            self.blocks.append(ParsedBlock(
                                block_type="title",
                                content=text,
                                page_num=slide_num,
                                section_path=list(self._section_stack),
                                source_file=self.ppt_path.name,
                            ))
                        else:
                            self.blocks.append(ParsedBlock(
                                block_type="text",
                                content=text,
                                page_num=slide_num,
                                section_path=list(self._section_stack),
                                source_file=self.ppt_path.name,
                            ))

            # ── 提取备注页 ──
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text and len(notes_text) > 5:
                    self.blocks.append(ParsedBlock(
                        block_type="text",
                        content=f"[备注] {notes_text}",
                        page_num=slide_num,
                        section_path=list(self._section_stack),
                        source_file=self.ppt_path.name,
                    ))

        logger.info(f"  PPT 解析完成: {len(self.blocks)} 个块")
        return self.blocks


# ── 统一入口 ──────────────────────────────────────────────────────────────────

def load_document(file_path: Path) -> list[ParsedBlock]:
    """
    根据文件后缀自动选择解析器。

    支持：
      .pdf  → PDFParser（pdfplumber + PyMuPDF）
      .pptx → PPTParser（python-pptx）
      .ppt  → PPTParser（python-pptx，兼容旧格式）
    """
    file_path = Path(file_path)
    if not file_path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    ext = file_path.suffix.lower()

    if ext == ".pdf":
        parser = PDFParser(file_path)
    elif ext in (".pptx", ".ppt"):
        parser = PPTParser(file_path)
    else:
        raise ValueError(f"不支持的文件格式: {ext}（仅支持 .pdf / .pptx / .ppt）")

    return parser.parse()


def save_parsed(blocks: list[ParsedBlock], source_path: Path):
    """将解析结果保存为 JSON，保留所有元信息。"""
    settings.parsed_dir.mkdir(parents=True, exist_ok=True)
    out_path = settings.parsed_dir / f"{source_path.stem}.json"

    output = {
        "source": str(source_path),
        "file_type": source_path.suffix.lower(),
        "block_count": len(blocks),
        "blocks": [asdict(b) for b in blocks],
    }
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(output, f, ensure_ascii=False, indent=2)
    logger.info(f"  已保存 → {out_path.name}（{len(blocks)} 个块）")
    return out_path


def load_all_documents() -> list[ParsedBlock]:
    """
    批量加载 data/raw_docs/ 下所有 PDF 和 PPT 文件。

    流程：
      1. 扫描 raw_docs 目录
      2. 逐个解析并保存到 parsed/ 目录
      3. 返回所有文档的 ParsedBlock 列表
    """
    settings.ensure_dirs()

    supported_exts = {".pdf", ".pptx", ".ppt"}
    files = [
        f for f in settings.raw_docs_dir.iterdir()
        if f.suffix.lower() in supported_exts
    ]

    if not files:
        logger.warning(
            f"未找到任何文档文件，请将 PDF/PPT 文件放入: {settings.raw_docs_dir}"
        )
        return []

    logger.info(f"发现 {len(files)} 个文档文件")

    all_blocks = []
    for file_path in sorted(files):
        try:
            blocks = load_document(file_path)
            save_parsed(blocks, file_path)
            all_blocks.extend(blocks)
        except Exception as e:
            logger.error(f"解析失败 {file_path.name}: {e}")

    logger.info(f"\n全部解析完成：共 {len(all_blocks)} 个块（来自 {len(files)} 个文件）")
    return all_blocks


# ── 命令行入口 ────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
    load_all_documents()
